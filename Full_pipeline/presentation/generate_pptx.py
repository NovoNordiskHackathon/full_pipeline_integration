#!/usr/bin/env python3
"""
Generate PowerPoint presentation for PTD Generator
Creates a basic .pptx file with slides and speaker notes
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

def create_presentation():
    """Create the PTD Generator presentation"""
    
    # Create presentation object
    prs = Presentation()
    
    # Define slide layouts
    title_slide_layout = prs.slide_layouts[0]  # Title slide
    content_slide_layout = prs.slide_layouts[1]  # Title and content
    two_content_layout = prs.slide_layouts[3]  # Two content
    
    # Slide 1: Title Slide
    slide1 = prs.slides.add_slide(title_slide_layout)
    title = slide1.shapes.title
    subtitle = slide1.placeholders[1]
    
    title.text = "PTD Generator"
    subtitle.text = "Automated Protocol Translation Document Generation\n\nTransform Protocol and CRF documents into structured PTD files using AI-powered processing"
    
    # Add speaker notes
    notes_slide = slide1.notes_slide
    text_frame = notes_slide.notes_text_frame
    text_frame.text = "Welcome to the PTD Generator presentation. This system automates the creation of Protocol Translation Documents from clinical trial documents. It combines AI-powered text extraction with structured data processing to streamline the clinical research workflow."
    
    # Slide 2: Project Overview
    slide2 = prs.slides.add_slide(content_slide_layout)
    title = slide2.shapes.title
    content = slide2.placeholders[1]
    
    title.text = "What is PTD Generator?"
    tf = content.text_frame
    tf.text = "Purpose: Transform Protocol PDFs and eCRF documents into structured PTD files"
    p = tf.add_paragraph()
    p.text = "Technology: AI-powered text extraction + hierarchical document structuring"
    p = tf.add_paragraph()
    p.text = "Output: Standardized Excel files with Schedule Grid and Study Specific Forms"
    p = tf.add_paragraph()
    p.text = "Target Users: Clinical research teams, data managers, protocol developers"
    
    # Add speaker notes
    notes_slide = slide2.notes_slide
    text_frame = notes_slide.notes_text_frame
    text_frame.text = "PTD Generator solves the manual, time-consuming process of creating Protocol Translation Documents. It uses advanced AI to extract and structure data from complex clinical documents, producing standardized outputs that integrate with existing clinical data management systems."
    
    # Slide 3: Architecture
    slide3 = prs.slides.add_slide(content_slide_layout)
    title = slide3.shapes.title
    content = slide3.placeholders[1]
    
    title.text = "System Architecture"
    tf = content.text_frame
    tf.text = "Three-Layer Architecture:"
    p = tf.add_paragraph()
    p.text = "• Frontend (HTML/JS) - User interface"
    p = tf.add_paragraph()
    p.text = "• Flask API (app.py) - Communication layer"
    p = tf.add_paragraph()
    p.text = "• PTD Pipeline (PTD_Gen/) - Processing modules"
    p = tf.add_paragraph()
    p.text = "• Modules: form_extractor, soa_parser, common_matrix, event_grouping, schedule_layout"
    
    # Add speaker notes
    notes_slide = slide3.notes_slide
    text_frame = notes_slide.notes_text_frame
    text_frame.text = "The system uses a three-layer architecture: Frontend, API, and Processing Pipeline. The Flask API serves as the communication layer between the web interface and processing modules. This modular design allows for easy maintenance and feature additions."
    
    # Slide 4: Project Structure
    slide4 = prs.slides.add_slide(content_slide_layout)
    title = slide4.shapes.title
    content = slide4.placeholders[1]
    
    title.text = "Project Structure"
    tf = content.text_frame
    tf.text = "Full_pipeline/"
    p = tf.add_paragraph()
    p.text = "├── app.py (Main Flask application)"
    p = tf.add_paragraph()
    p.text = "├── frontend/ (Web interface)"
    p = tf.add_paragraph()
    p.text = "├── PTD_Gen/ (Core PTD generation)"
    p = tf.add_paragraph()
    p.text = "└── docs/ (Documentation)"
    
    # Add speaker notes
    notes_slide = slide4.notes_slide
    text_frame = notes_slide.notes_text_frame
    text_frame.text = "The project has a clean, organized structure with clear separation of concerns. The frontend handles user interaction, PTD_Gen contains core processing logic, and comprehensive documentation ensures maintainability."
    
    # Slide 5: Key Modules
    slide5 = prs.slides.add_slide(two_content_layout)
    title = slide5.shapes.title
    left_content = slide5.placeholders[1]
    right_content = slide5.placeholders[2]
    
    title.text = "Key Modules"
    
    # Left side - Form Extraction
    tf_left = left_content.text_frame
    tf_left.text = "Form Extraction Module"
    p = tf_left.add_paragraph()
    p.text = "• Extract forms from eCRF JSON"
    p = tf_left.add_paragraph()
    p.text = "• Classify form sources"
    p = tf_left.add_paragraph()
    p.text = "• Identify triggers & requirements"
    p = tf_left.add_paragraph()
    p.text = "• Generate structured CSV"
    
    # Right side - Schedule Parser
    tf_right = right_content.text_frame
    tf_right.text = "Schedule Parser Module"
    p = tf_right.add_paragraph()
    p.text = "• Parse schedule tables"
    p = tf_right.add_paragraph()
    p.text = "• Extract visit patterns"
    p = tf_right.add_paragraph()
    p.text = "• Create visit-procedure matrix"
    p = tf_right.add_paragraph()
    p.text = "• Handle complex structures"
    
    # Add speaker notes
    notes_slide = slide5.notes_slide
    text_frame = notes_slide.notes_text_frame
    text_frame.text = "The form extraction module uses pattern matching and AI to intelligently extract form information. The schedule parser handles various visit naming conventions and provides robust parsing for different document formats."
    
    # Slide 6: API Endpoints
    slide6 = prs.slides.add_slide(content_slide_layout)
    title = slide6.shapes.title
    content = slide6.placeholders[1]
    
    title.text = "RESTful API Interface"
    tf = content.text_frame
    tf.text = "Key Endpoints:"
    p = tf.add_paragraph()
    p.text = "• GET /status - System status and health"
    p = tf.add_paragraph()
    p.text = "• POST /run_pipeline - Main processing endpoint"
    p = tf.add_paragraph()
    p.text = "• POST /run_ptd_generation - Full PTD generation"
    p = tf.add_paragraph()
    p.text = "• GET /download/<filename> - File download"
    p = tf.add_paragraph()
    p.text = "• Supports both file upload and JSON data input"
    
    # Add speaker notes
    notes_slide = slide6.notes_slide
    text_frame = notes_slide.notes_text_frame
    text_frame.text = "The system provides a RESTful API design for easy integration with existing systems. It supports both file upload and JSON data input, with comprehensive error handling and status reporting."
    
    # Slide 7: Setup & Run
    slide7 = prs.slides.add_slide(content_slide_layout)
    title = slide7.shapes.title
    content = slide7.placeholders[1]
    
    title.text = "Quick Setup and Running"
    tf = content.text_frame
    tf.text = "Setup Commands:"
    p = tf.add_paragraph()
    p.text = "python -m venv venv"
    p = tf.add_paragraph()
    p.text = "source venv/bin/activate"
    p = tf.add_paragraph()
    p.text = "pip install -r requirements.txt"
    p = tf.add_paragraph()
    p.text = "python app.py"
    p = tf.add_paragraph()
    p.text = "Access: http://localhost:5000"
    
    # Add speaker notes
    notes_slide = slide7.notes_slide
    text_frame = notes_slide.notes_text_frame
    text_frame.text = "The setup process is simple with minimal dependencies. The application runs on standard port 5000 and is easy to deploy and maintain."
    
    # Slide 8: Troubleshooting
    slide8 = prs.slides.add_slide(content_slide_layout)
    title = slide8.shapes.title
    content = slide8.placeholders[1]
    
    title.text = "Common Issues and Solutions"
    tf = content.text_frame
    tf.text = "Top 5 Fixes:"
    p = tf.add_paragraph()
    p.text = "1. ModuleNotFoundError: Add current directory to PYTHONPATH"
    p = tf.add_paragraph()
    p.text = "2. Port Already in Use: Kill existing process or use different port"
    p = tf.add_paragraph()
    p.text = "3. Missing Dependencies: Reinstall requirements.txt"
    p = tf.add_paragraph()
    p.text = "4. File Not Found: Check file paths and permissions"
    p = tf.add_paragraph()
    p.text = "5. JSON Decode Error: Validate input JSON files"
    
    # Add speaker notes
    notes_slide = slide8.notes_slide
    text_frame = notes_slide.notes_text_frame
    text_frame.text = "Common issues are well-documented with specific solutions. Debug commands help identify and resolve problems quickly, and comprehensive logging provides detailed information for troubleshooting."
    
    # Slide 9: Demo Flow
    slide9 = prs.slides.add_slide(content_slide_layout)
    title = slide9.shapes.title
    content = slide9.placeholders[1]
    
    title.text = "Live Demo - PTD Generation"
    tf = content.text_frame
    tf.text = "Demo Steps:"
    p = tf.add_paragraph()
    p.text = "1. Start Application: python app.py"
    p = tf.add_paragraph()
    p.text = "2. Access Frontend: http://localhost:5000"
    p = tf.add_paragraph()
    p.text = "3. Upload Files: Drag and drop Protocol and CRF documents"
    p = tf.add_paragraph()
    p.text = "4. Process Documents: Click 'Generate PTD' button"
    p = tf.add_paragraph()
    p.text = "5. Monitor Progress: Watch processing steps"
    p = tf.add_paragraph()
    p.text = "6. Download Result: Get generated PTD Excel file"
    
    # Add speaker notes
    notes_slide = slide9.notes_slide
    text_frame = notes_slide.notes_text_frame
    text_frame.text = "This live demonstration shows the complete workflow from document upload to PTD generation. The user-friendly interface and processing capabilities are demonstrated, along with real-world usage scenarios."
    
    # Slide 10: Next Steps
    slide10 = prs.slides.add_slide(content_slide_layout)
    title = slide10.shapes.title
    content = slide10.placeholders[1]
    
    title.text = "Future Enhancements"
    tf = content.text_frame
    tf.text = "Immediate Improvements:"
    p = tf.add_paragraph()
    p.text = "• Complete integration testing"
    p = tf.add_paragraph()
    p.text = "• Enhanced error handling"
    p = tf.add_paragraph()
    p.text = "• Input validation"
    p = tf.add_paragraph()
    p.text = "Short-term Goals:"
    p = tf.add_paragraph()
    p.text = "• Docker containerization"
    p = tf.add_paragraph()
    p.text = "• CI/CD pipeline"
    p = tf.add_paragraph()
    p.text = "• API documentation with Swagger"
    p = tf.add_paragraph()
    p.text = "Long-term Vision:"
    p = tf.add_paragraph()
    p.text = "• Microservices architecture"
    p = tf.add_paragraph()
    p.text = "• Cloud deployment"
    p = tf.add_paragraph()
    p.text = "• Machine learning enhancements"
    
    # Add speaker notes
    notes_slide = slide10.notes_slide
    text_frame = notes_slide.notes_text_frame
    text_frame.text = "The roadmap includes clear priorities for future development, focusing on scalability and maintainability. The community-driven development approach ensures continuous improvement and adaptation to user needs."
    
    # Slide 11: Documentation
    slide11 = prs.slides.add_slide(content_slide_layout)
    title = slide11.shapes.title
    content = slide11.placeholders[1]
    
    title.text = "Documentation and Support"
    tf = content.text_frame
    tf.text = "Documentation:"
    p = tf.add_paragraph()
    p.text = "• docs/README_FULL.md - Comprehensive documentation"
    p = tf.add_paragraph()
    p.text = "• docs/ARCHITECTURE.md - System architecture"
    p = tf.add_paragraph()
    p.text = "• docs/API_REFERENCE.md - API documentation"
    p = tf.add_paragraph()
    p.text = "• docs/QUICK_START.md - Quick start guide"
    p = tf.add_paragraph()
    p.text = "Support:"
    p = tf.add_paragraph()
    p.text = "• GitHub repository with issue tracking"
    p = tf.add_paragraph()
    p.text = "• Comprehensive troubleshooting guide"
    p = tf.add_paragraph()
    p.text = "• Example configurations and templates"
    
    # Add speaker notes
    notes_slide = slide11.notes_slide
    text_frame = notes_slide.notes_text_frame
    text_frame.text = "Extensive documentation is available for both users and developers, with multiple support channels for assistance. Regular updates and community contributions ensure the project remains current and responsive to user needs."
    
    # Slide 12: Questions & Discussion
    slide12 = prs.slides.add_slide(content_slide_layout)
    title = slide12.shapes.title
    content = slide12.placeholders[1]
    
    title.text = "Questions and Discussion"
    tf = content.text_frame
    tf.text = "Contact Information:"
    p = tf.add_paragraph()
    p.text = "• Repository: [GitHub Link]"
    p = tf.add_paragraph()
    p.text = "• Documentation: docs/ directory"
    p = tf.add_paragraph()
    p.text = "• Issues: GitHub Issues tracker"
    p = tf.add_paragraph()
    p.text = "Discussion Topics:"
    p = tf.add_paragraph()
    p.text = "• Integration requirements"
    p = tf.add_paragraph()
    p.text = "• Customization needs"
    p = tf.add_paragraph()
    p.text = "• Performance considerations"
    p = tf.add_paragraph()
    p.text = "• Deployment strategies"
    
    # Add speaker notes
    notes_slide = slide12.notes_slide
    text_frame = notes_slide.notes_text_frame
    text_frame.text = "This is an open floor for questions and discussion. We can address specific use cases and requirements, gather feedback for future improvements, and discuss integration strategies."
    
    return prs

def main():
    """Generate the PowerPoint presentation"""
    try:
        # Create presentation
        prs = create_presentation()
        
        # Save presentation
        output_path = "/workspace/Full_pipeline/presentation/Full_Pipeline_Presentation.pptx"
        prs.save(output_path)
        
        print(f"✅ Presentation created successfully: {output_path}")
        print(f"📊 Total slides: {len(prs.slides)}")
        print("📝 Speaker notes included for each slide")
        
    except ImportError:
        print("❌ Error: python-pptx library not installed")
        print("💡 Install with: pip install python-pptx")
    except Exception as e:
        print(f"❌ Error creating presentation: {e}")

if __name__ == "__main__":
    main()